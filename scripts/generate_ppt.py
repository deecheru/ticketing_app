import os
import django
import sys
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
from django.db.models import Count

# Set up Django environment
sys.path.append(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
os.environ.setdefault('DJANGO_SETTINGS_MODULE', 'ticketing_system.settings')
django.setup()

# Import models after Django setup
from user_tickets.models import Ticket, ServiceFamily, ServiceType, ServiceCategory
from accounts.models import Company, User

def create_title_slide(prs):
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = title_slide.shapes.title
    subtitle = title_slide.placeholders[1]
    
    title.text = "Ticketing System Overview"
    subtitle.text = "Enterprise Support Management Platform"

def create_statistics_slide(prs):
    bullet_slide = prs.slides.add_slide(prs.slide_layouts[1])
    shapes = bullet_slide.shapes
    
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = 'System Statistics'
    
    tf = body_shape.text_frame
    tf.text = f"Total Tickets: {Ticket.objects.count()}"
    
    p = tf.add_paragraph()
    p.text = f"Total Companies: {Company.objects.count()}"
    
    p = tf.add_paragraph()
    p.text = f"Active Users: {User.objects.filter(is_active=True).count()}"
    
    p = tf.add_paragraph()
    p.text = f"Service Families: {ServiceFamily.objects.count()}"

def create_ticket_status_slide(prs):
    chart_slide = prs.slides.add_slide(prs.slide_layouts[6])
    shapes = chart_slide.shapes
    
    title_shape = shapes.title
    title_shape.text = 'Ticket Status Distribution'
    
    # Get ticket status counts
    status_counts = Ticket.objects.values('status').annotate(count=Count('id'))
    
    # Add text box with statistics
    left = Inches(1)
    top = Inches(2)
    for status in status_counts:
        txBox = shapes.add_textbox(left, top, Inches(6), Inches(0.5))
        tf = txBox.text_frame
        tf.text = f"{status['status']}: {status['count']} tickets"
        top += Inches(0.6)

def generate_presentation():
    prs = Presentation()
    
    # Create slides
    create_title_slide(prs)
    create_statistics_slide(prs)
    create_ticket_status_slide(prs)
    
    # Save presentation
    output_path = os.path.join(os.path.dirname(os.path.dirname(__file__)), 
                              'presentations', 
                              'system_overview.pptx')
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")

if __name__ == "__main__":
    generate_presentation()