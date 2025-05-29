from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_title_slide(prs):
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Ticketing System"
    subtitle.text = "Project Presentation"
    
    # Format title
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

def create_section_slide(prs, title, content):
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = title
    
    tf = body_shape.text_frame
    for item in content:
        p = tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(18)

def create_presentation():
    prs = Presentation()
    
    # Title Slide
    create_title_slide(prs)
    
    # System Architecture
    architecture_content = [
        "Built with Django (Python)",
        "Frontend: HTML, CSS, JavaScript, Bootstrap",
        "Database: Relational database",
        "Key Components:",
        "• Accounts Management",
        "• Ticket Management",
        "• Service Hierarchy",
        "• File Storage System"
    ]
    create_section_slide(prs, "System Architecture", architecture_content)
    
    # Key Features
    features_content = [
        "Multi-tenant Architecture",
        "Role-based Access Control",
        "Comprehensive Ticket Management",
        "Service Hierarchy System",
        "File Attachment Support",
        "Notification System",
        "Reporting and Analytics"
    ]
    create_section_slide(prs, "Key Features", features_content)
    
    # User Roles
    roles_content = [
        "Superadmin (System Administrator)",
        "• Complete system access",
        "• Manage all companies and users",
        "Company Agent",
        "• Manage tickets from assigned companies",
        "Company Staff",
        "• Manage company-specific tickets",
        "Regular User",
        "• Create and view own tickets"
    ]
    create_section_slide(prs, "User Roles", roles_content)
    
    # Ticket Lifecycle
    lifecycle_content = [
        "Creation",
        "• User submits ticket with details",
        "Processing",
        "• Staff reviews and updates ticket",
        "Resolution",
        "• Issue is fixed and documented",
        "Closure",
        "• Ticket is verified and closed",
        "• Archived for future reference"
    ]
    create_section_slide(prs, "Ticket Lifecycle", lifecycle_content)
    
    # Service Hierarchy
    hierarchy_content = [
        "Three-level Categorization:",
        "Service Family",
        "• Top-level categorization",
        "Service Type",
        "• Mid-level categorization",
        "Service Category",
        "• Detailed categorization",
        "Benefits:",
        "• Organized ticket routing",
        "• Accurate tracking",
        "• Detailed reporting"
    ]
    create_section_slide(prs, "Service Hierarchy", hierarchy_content)
    
    # Security Features
    security_content = [
        "Role-based Access Control",
        "Multi-factor Authentication",
        "Secure File Storage",
        "Company Data Isolation",
        "Permission Management",
        "Regular Security Audits"
    ]
    create_section_slide(prs, "Security Features", security_content)
    
    # Benefits for Clients
    benefits_content = [
        "Streamlined Ticket Management",
        "Improved Response Times",
        "Better Organization",
        "Enhanced Security",
        "Scalable Solution",
        "Customizable Service Categories",
        "Comprehensive Reporting"
    ]
    create_section_slide(prs, "Benefits for Clients", benefits_content)
    
    # Save the presentation
    prs.save('ticketing_system_presentation.pptx')

if __name__ == "__main__":
    create_presentation() 